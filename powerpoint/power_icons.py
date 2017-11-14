'''
Script to create powerpoint with icons from the emf files.

Author: Dave Vieglais
'''

import glob
import logging
from pptx import Presentation
from pptx.util import Inches

# Index to a blank page template in the PPTX
BLANK_LAYOUT_INDEX = 3


def addIconPage(ppt, icons, title, start_index=0, rows=5, cols=10, TLx=1.0, TLy=1.5, width=8.0, height=3.5):
  '''
  Adds a matrix of icons to a new page

  :param ppt: instance of Powerpoint
  :param icons: list of paths to image files
  :param title: Title for the generated page
  :param start_index: zero-based start index into icons
  :param rows: number of rows on the page
  :param cols: number of columns on the page
  :param TLx: top left X position, inches
  :param TLy: top left Y postion from top, inches
  :param width: width in inches of writable area
  :param height: height in index of writable area
  :return: offset in icons of next icon
  '''
  blank_slide_layout = ppt.slide_layouts[ BLANK_LAYOUT_INDEX ]
  page = ppt.slides.add_slide(blank_slide_layout)
  page.shapes.title.text = title
  dx = width/(cols*1.0)
  dy = height/(rows*1.0)
  ico_height = Inches(dy*0.80)
  i = start_index
  for row in range(0,rows):
    y = Inches(TLy + dy*row)
    for col in range(0,cols):
      x = Inches(TLx + dx*col)
      logging.debug("Adding %s at (%.2f, %.2f)", icons[i], x, y)
      page.shapes.add_picture(icons[i], x, y, height=ico_height)
      i += 1
      if i >= len(icons):
        return i
  return i


def addETLine(ppt, path="et-line/emf/*.emf"):
  '''
  Add the et-line icon pages
  :param ppt: instance of Powerpoint
  :param path: path to search for emf files
  :return: nothin
  '''
  icon_files = glob.glob(path, recursive=False)
  icon_files.sort()
  idx = 0
  page = 1
  TITLE = "ET-Line {}"
  while idx < len(icon_files):
    logging.info("Processing et-line from index=%d", idx)
    title = TITLE.format(page)
    idx = addIconPage(ppt, icon_files, title, start_index=idx)
    page += 1


def addFontAwesome(ppt, path="font-awesome/emf/*.emf"):
  '''
  Add the font-awesome icon pages
  :param ppt: instance of Powerpoint
  :param path: path to search for emf files
  :return: nothin
  '''
  icon_files = glob.glob(path, recursive=False)
  icon_files.sort()
  idx = 0
  page = 1
  TITLE = "Font-Awesome {}"
  while idx < len(icon_files):
    logging.info("Processing font-awesome from index=%d", idx)
    title = TITLE.format(page)
    idx = addIconPage(ppt, icon_files, title, start_index=idx)
    page += 1


if __name__ == "__main__":
  logging.basicConfig(level=logging.DEBUG)
  ppt = Presentation("template.pptx")
  addETLine(ppt)
  addFontAwesome(ppt)
  ppt.save("icons.pptx")
